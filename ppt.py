from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Smart Lock Market Analysis"
subtitle.text = "Circuit House Assignment"

# Analysis slides
# Slide for number of brands
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Number of Brands"
body.text = f"There are {num_brands} brands in the smart lock segment."

# Slide for SKUs per brand
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "SKUs per Brand"
body.text = "\n".join([f"{brand}: {count}" for brand, count in skus_per_brand.items()])

# Slide for relative ranking
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Relative Ranking"
body.text = "\n".join([f"{brand}: {rank:.2f}" for brand, rank in relative_ranks.items()])

# Slide for relative ratings
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Relative Ratings"
body.text = "\n".join([f"{brand}: {rating:.2f}" for brand, rating in relative_ratings.items()])

# Slide for price distribution
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Price Distribution"
body.text = "\n".join([f"{band}: {count}" for band, count in price_distribution.items()])

# Save the presentation
prs.save('smart_lock_analysis.pptx')